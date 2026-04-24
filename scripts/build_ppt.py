#!/usr/bin/env python3
"""
Build the combined presentation:
  Slide 1  : Title
  Slide 2  : Problem Statement
  Slides 3-7  : Our EMBER2024 model (5 slides, with actual plot images)
  Slide 8  : Section divider — Friend's model
  Slides 9-13 : Friend's Random Forest APK model (5 slides)
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree
import copy

FIGS     = Path("results/figs")
FIGS_RF  = Path("results/figs/rf_friend")
OUT      = Path("results/EMBER2024_presentation.pptx")

# ── colour palette ─────────────────────────────────────────────────────────────
BG_DARK   = RGBColor(0x0D, 0x1B, 0x2A)   # deep navy
BG_FRIEND = RGBColor(0x0A, 0x27, 0x1A)   # deep green-navy
ACCENT1   = RGBColor(0x1F, 0x77, 0xB4)   # blue
ACCENT2   = RGBColor(0xFF, 0x7F, 0x0E)   # orange
ACCENT3   = RGBColor(0x2C, 0xA0, 0x2C)   # green
ACCENT_R  = RGBColor(0xD6, 0x27, 0x28)   # red
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT     = RGBColor(0xCC, 0xDD, 0xEE)
YELLOW    = RGBColor(0xFF, 0xD7, 0x00)
GREY      = RGBColor(0xAA, 0xBB, 0xCC)

W = Inches(13.33)   # widescreen 16:9
H = Inches(7.5)


# ── helpers ────────────────────────────────────────────────────────────────────
def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(layout)


def solid_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill: RGBColor, alpha=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE_TYPE.RECTANGLE
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h,
             size=20, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb


def add_bullets(slide, items, x, y, w, h,
                size=16, color=WHITE, title_color=ACCENT1,
                marker="▸ ", line_spacing=1.15):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(3)
        run = p.add_run()
        is_header = item.startswith("##")
        is_sub    = item.startswith("  ")
        txt = item.lstrip("# ").lstrip()
        run.text = ("" if is_sub else marker) + txt
        run.font.size  = Pt(size - 2 if is_sub else size)
        run.font.bold  = is_header
        run.font.color.rgb = title_color if is_header else (GREY if is_sub else color)
    return txb


def add_image(slide, path, x, y, w, h=None):
    if not Path(path).exists():
        return None
    if h:
        slide.shapes.add_picture(str(path), x, y, w, h)
    else:
        slide.shapes.add_picture(str(path), x, y, w)


def divider_line(slide, y, color=ACCENT1):
    line = slide.shapes.add_shape(1, Inches(0.5), y, Inches(12.33), Pt(2))
    line.fill.solid(); line.fill.fore_color.rgb = color
    line.line.fill.background()


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ══════════════════════════════════════════════════════════════════════════════

def slide_title(prs):
    s = blank_slide(prs); solid_bg(s, BG_DARK)
    add_rect(s, 0, Inches(2.8), W, Inches(2.0), RGBColor(0x1F, 0x77, 0xB4))
    add_text(s, "EMBER2024 Malware Detection", Inches(0.6), Inches(3.0),
             Inches(12), Inches(1.0), size=40, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, "Beating the Paper Baseline with LightGBM Ensembling  +  Deep Learning Pipeline",
             Inches(0.6), Inches(3.9), Inches(12), Inches(0.6),
             size=18, color=LIGHT, align=PP_ALIGN.CENTER)
    add_text(s, "DACS Project  ·  2024", Inches(0.6), Inches(5.2),
             Inches(12), Inches(0.5), size=14, color=GREY, align=PP_ALIGN.CENTER)


def slide_problem(prs):
    s = blank_slide(prs); solid_bg(s, BG_DARK)
    add_rect(s, 0, 0, W, Inches(1.1), RGBColor(0x0A, 0x3D, 0x6B))
    add_text(s, "Problem Statement", Inches(0.5), Inches(0.15),
             Inches(12), Inches(0.8), size=30, bold=True)
    divider_line(s, Inches(1.1))

    left_items = [
        "##The Challenge",
        "Malware evolves faster than static signatures",
        "Evasive malware is crafted to bypass detection",
        "6,315 challenge samples — designed to look benign",
        "##Two Tasks",
        "Binary detection  (benign vs malicious)",
        "Family classification  (long-tail distribution)",
        "##Why It's Hard",
        "Test ROC-AUC is saturated (~0.997) — challenge set separates models",
        "Score scale differs per file type (PDF ≠ Win32)",
    ]
    right_items = [
        "##Dataset — EMBER2024",
        "2,626,000 training samples",
        "605,929 test samples  (balanced 50/50)",
        "6,315 evasive challenge samples (all malicious)",
        "6 file types: Win32, Win64, .NET, APK, ELF, PDF",
        "2,568-dim feature vector per file (static features)",
        "##Evaluation Protocol",
        "Challenge ROC: evasive malware vs benign test samples",
        "Detection rate = fraction of challenge samples scored ≥ 0.5",
        "Paper baseline: EMBER2024_all.model (500 trees, 64 leaves)",
    ]
    add_bullets(s, left_items,  Inches(0.5), Inches(1.2), Inches(6.0), Inches(5.8), size=15)
    add_bullets(s, right_items, Inches(7.0), Inches(1.2), Inches(6.0), Inches(5.8), size=15)
    add_rect(s, Inches(6.75), Inches(1.2), Pt(2), Inches(5.8), ACCENT1)


def slide_our_approach(prs):
    s = blank_slide(prs); solid_bg(s, BG_DARK)
    add_rect(s, 0, 0, W, Inches(1.1), RGBColor(0x0A, 0x3D, 0x6B))
    add_text(s, "Our Approach — Two-Track Pipeline", Inches(0.5), Inches(0.15),
             Inches(12), Inches(0.8), size=28, bold=True)
    divider_line(s, Inches(1.1))

    add_image(s, FIGS/"01_architecture.png",
              Inches(0.3), Inches(1.15), Inches(12.7), Inches(5.4))


def slide_our_metrics(prs):
    s = blank_slide(prs); solid_bg(s, BG_DARK)
    add_rect(s, 0, 0, W, Inches(1.1), RGBColor(0x0A, 0x3D, 0x6B))
    add_text(s, "Evaluation Metrics — All Models vs Paper", Inches(0.5), Inches(0.15),
             Inches(12), Inches(0.8), size=28, bold=True)
    divider_line(s, Inches(1.1))

    add_image(s, FIGS/"02_metrics_table.png",
              Inches(0.5), Inches(1.2), Inches(12.3), Inches(3.0))

    steps = [
        "##Step-by-step improvements",
        "2048-leaf model  →  47.70% det. rate  (overfits)",
        "Paper baseline   →  66.54%  (reference)",
        "64-leaf retrain  →  65.57%  (better generalisation)",
        "Per-type paper models  →  70.28%",
        "Grid-search weights (V2)  →  70.23%, PR-AUC 0.6284  ✓ beats paper all metrics",
        "Rank ensemble (novel)  →  97.86% det. rate  ✓ +31pp over paper",
    ]
    add_bullets(s, steps, Inches(0.5), Inches(4.3), Inches(12.3), Inches(3.0), size=14)


def slide_our_roc(prs):
    s = blank_slide(prs); solid_bg(s, BG_DARK)
    add_rect(s, 0, 0, W, Inches(1.1), RGBColor(0x0A, 0x3D, 0x6B))
    add_text(s, "ROC & Precision-Recall Curves", Inches(0.5), Inches(0.15),
             Inches(12), Inches(0.8), size=28, bold=True)
    divider_line(s, Inches(1.1))

    add_image(s, FIGS/"03_roc_curves.png",
              Inches(0.2), Inches(1.15), Inches(13.0), Inches(3.05))
    add_image(s, FIGS/"04_pr_curves.png",
              Inches(0.2), Inches(4.2),  Inches(13.0), Inches(3.05))


def slide_our_features(prs):
    s = blank_slide(prs); solid_bg(s, BG_DARK)
    add_rect(s, 0, 0, W, Inches(1.1), RGBColor(0x0A, 0x3D, 0x6B))
    add_text(s, "Feature Importance — Our 64-leaf LightGBM", Inches(0.5), Inches(0.15),
             Inches(12), Inches(0.8), size=28, bold=True)
    divider_line(s, Inches(1.1))

    add_image(s, FIGS/"05_feature_importance.png",
              Inches(0.3), Inches(1.15), Inches(8.5), Inches(5.5))

    bullets = [
        "##Key findings",
        "Byte histogram features (idx 256–1279) dominate by gain",
        "Header features (idx 0–7) dominate by split count",
        "String statistics & import table are highly informative",
        "ELF/APK-specific features rarely used by Win32-trained models",
        "##Why per-type models help",
        "Win32 model ignores APK-specific dimensions (zero gain)",
        "Per-type training = feature selection for free",
        "2,568 features → per-type model uses ~300 effectively",
    ]
    add_bullets(s, bullets, Inches(9.0), Inches(1.3), Inches(4.1), Inches(5.3), size=14)


def slide_our_failure(prs):
    s = blank_slide(prs); solid_bg(s, BG_DARK)
    add_rect(s, 0, 0, W, Inches(1.1), RGBColor(0x0A, 0x3D, 0x6B))
    add_text(s, "Ensemble Results & Failure Analysis", Inches(0.5), Inches(0.15),
             Inches(12), Inches(0.8), size=28, bold=True)
    divider_line(s, Inches(1.1))

    add_image(s, FIGS/"06_ensemble_failure.png",
              Inches(0.2), Inches(1.15), Inches(13.0), Inches(4.3))

    bullets = [
        "##Novel contribution: Rank-based ensemble",
        "Paper uses raw score averaging — scale bias across file types",
        "Borda count ranks + power sharpening (T=2) removes scale bias",
        "PDF model's 0.7 ≠ Win32 model's 0.7 — ranking normalises this",
        "Trade-off: detection rate 97.86% vs ROC-AUC 0.9496 (slightly below paper)",
        "V2 raw ensemble beats paper on ALL 4 metrics simultaneously",
    ]
    add_bullets(s, bullets, Inches(0.5), Inches(5.5), Inches(12.3), Inches(1.8), size=13)


# ── section divider ────────────────────────────────────────────────────────────
def slide_divider(prs, title, subtitle, bg=RGBColor(0x0A, 0x27, 0x1A), accent=ACCENT3):
    s = blank_slide(prs); solid_bg(s, bg)
    add_rect(s, 0, Inches(2.9), W, Inches(1.8), accent)
    add_text(s, title,    Inches(0.6), Inches(3.0), Inches(12), Inches(1.0),
             size=38, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, subtitle, Inches(0.6), Inches(3.9), Inches(12), Inches(0.6),
             size=18, color=LIGHT, align=PP_ALIGN.CENTER)


# ── friend's slides ────────────────────────────────────────────────────────────
def slide_rf_overview(prs):
    s = blank_slide(prs); solid_bg(s, BG_DARK)
    add_rect(s, 0, 0, W, Inches(1.1), RGBColor(0x0A, 0x3D, 0x6B))
    add_text(s, "Problem & Model Overview — Random Forest (PDF)", Inches(0.5), Inches(0.15),
             Inches(12.5), Inches(0.8), size=26, bold=True, color=WHITE)
    divider_line(s, Inches(1.1), ACCENT1)

    left = [
        "##Problem",
        "Binary malware classification on PDF files",
        "Real-world challenge: PDF malware uses embedded scripts & exploits",
        "Goal: high precision + recall on unseen PDF samples",
        "##Dataset",
        "24,000 PDF test samples from EMBER2024",
        "Static features extracted from PDF structure & content",
        "Class-imbalanced — malicious PDFs minority class",
    ]
    right = [
        "##Model: Random Forest (chunked training)",
        "Trained in chunks to handle memory constraints",
        "Ensemble of decision trees — robust to feature noise",
        "No gradient information needed — fast, interpretable",
        "##Why Random Forest?",
        "Handles high-dimensional sparse features well",
        "Naturally provides feature importance (Gini impurity)",
        "Resistant to overfitting on noisy PDF static features",
        "Out-of-bag error estimation built-in",
    ]
    add_bullets(s, left,  Inches(0.5), Inches(1.2), Inches(5.9), Inches(5.8), size=15, title_color=ACCENT1)
    add_bullets(s, right, Inches(6.9), Inches(1.2), Inches(6.0), Inches(5.8), size=15, title_color=ACCENT1)
    add_rect(s, Inches(6.65), Inches(1.2), Pt(2), Inches(5.8), ACCENT1)


def slide_rf_metrics(prs):
    s = blank_slide(prs); solid_bg(s, BG_DARK)
    add_rect(s, 0, 0, W, Inches(1.1), RGBColor(0x0A, 0x3D, 0x6B))
    add_text(s, "Model Performance Metrics  —  PDF (n=24,000)", Inches(0.5), Inches(0.15),
             Inches(12.5), Inches(0.8), size=26, bold=True, color=WHITE)
    divider_line(s, Inches(1.1), ACCENT1)

    # Big metric boxes — actual values from rf_metrics.json
    metrics = [
        ("ROC-AUC",   "0.9888", ACCENT1),
        ("PR-AUC",    "0.9909", ACCENT3),
        ("Precision", "0.9945", ACCENT2),
        ("Recall",    "0.8992", RGBColor(0x9B, 0x59, 0xB6)),
        ("F1-Score",  "0.9444", ACCENT_R),
    ]
    box_w, box_h = Inches(2.2), Inches(1.5)
    gap = Inches(0.3)
    start_x = Inches(0.65)
    y_box = Inches(1.4)
    for i, (name, val, col) in enumerate(metrics):
        x = start_x + i * (box_w + gap)
        add_rect(s, x, y_box, box_w, box_h, col)
        add_text(s, val,  x, y_box + Inches(0.1), box_w, Inches(0.9),
                 size=30, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, name, x, y_box + Inches(0.95), box_w, Inches(0.45),
                 size=13, align=PP_ALIGN.CENTER, color=WHITE)

    bullets = [
        "##Interpretation",
        "ROC-AUC 0.9888 — excellent discrimination across all thresholds on PDF files",
        "PR-AUC 0.9909 — very high precision maintained even at high recall",
        "Precision 0.9945 — only 0.55% of flagged PDFs are false positives",
        "Recall 0.8992 — catches ~90% of malicious PDFs; ~10% evasive samples missed",
        "F1-Score 0.9444 — strong harmonic balance of precision & recall",
        "##Comparison context",
        "Strong result on a hard file type — PDF malware uses embedded JS, exploits, obfuscation",
        "Precision far exceeds recall — tuned for low false alarm rate in production",
    ]
    add_bullets(s, bullets, Inches(0.5), Inches(3.1), Inches(12.3), Inches(4.0),
                size=14, title_color=ACCENT1)


def slide_rf_curves(prs):
    s = blank_slide(prs); solid_bg(s, BG_DARK)
    add_rect(s, 0, 0, W, Inches(1.1), RGBColor(0x0A, 0x3D, 0x6B))
    add_text(s, "ROC & Precision-Recall Analysis", Inches(0.5), Inches(0.15),
             Inches(12.5), Inches(0.8), size=28, bold=True, color=WHITE)
    divider_line(s, Inches(1.1), ACCENT1)

    # Left: actual ROC plot
    add_image(s, FIGS_RF / "rf_roc.png",
              Inches(0.3), Inches(1.2), Inches(6.3), Inches(4.8))

    # Right: actual PR plot
    add_image(s, FIGS_RF / "rf_pr.png",
              Inches(6.8), Inches(1.2), Inches(6.3), Inches(4.8))

    notes = [
        "ROC-AUC 0.9888: curve hugs top-left corner — very low FPR at high TPR",
        "PR-AUC 0.9909: precision stays near 1.0 until recall ≈ 0.85 — ideal for low-FP production use",
    ]
    add_bullets(s, notes, Inches(0.4), Inches(6.1), Inches(12.5), Inches(1.2), size=13,
                title_color=ACCENT1)


def slide_rf_features(prs):
    s = blank_slide(prs); solid_bg(s, BG_DARK)
    add_rect(s, 0, 0, W, Inches(1.1), RGBColor(0x0A, 0x3D, 0x6B))
    add_text(s, "Feature Importance Insights — PDF", Inches(0.5), Inches(0.15),
             Inches(12.5), Inches(0.8), size=28, bold=True, color=WHITE)
    divider_line(s, Inches(1.1), ACCENT1)

    # Actual feature importance plot
    add_image(s, FIGS_RF / "rf_features.png",
              Inches(0.3), Inches(1.2), Inches(8.7), Inches(5.5))

    bullets = [
        "##Top PDF features",
        "Feature 622 dominates (importance ≈ 0.044)",
        "Byte entropy features — core signal for PDF malware",
        "PDF object-level structure features",
        "Header/metadata fields indicate obfuscation",
        "##Why feature importance matters",
        "Reveals PDF-specific malware patterns",
        "Byte entropy = high randomness → encrypted payloads",
        "Guides feature engineering & dimensionality reduction",
        "Provides explainability for security analysts",
    ]
    add_bullets(s, bullets, Inches(9.2), Inches(1.3), Inches(3.9), Inches(5.5),
                size=13, title_color=ACCENT1)


def slide_rf_takeaways(prs):
    s = blank_slide(prs); solid_bg(s, BG_DARK)
    add_rect(s, 0, 0, W, Inches(1.1), RGBColor(0x0A, 0x3D, 0x6B))
    add_text(s, "Key Takeaways & Issues", Inches(0.5), Inches(0.15),
             Inches(12.5), Inches(0.8), size=28, bold=True, color=WHITE)
    divider_line(s, Inches(1.1), ACCENT1)

    # Green: strengths
    add_rect(s, Inches(0.4), Inches(1.25), Inches(5.9), Inches(3.0), RGBColor(0x0A, 0x2A, 0x45))
    add_text(s, "✓  Random Forest — Strengths", Inches(0.5), Inches(1.3),
             Inches(5.7), Inches(0.5), size=15, bold=True, color=ACCENT1)
    rf_good = [
        "ROC-AUC 0.9888, PR-AUC 0.9909 — excellent discrimination on PDF files",
        "Precision 0.9945 — only 0.55% false alarm rate (very few false positives)",
        "Robust to PDF feature noise from obfuscated/encrypted payloads",
        "Fast inference — chunked training handles large-scale PDF telemetry",
        "Gini importance reveals key features for analyst explainability",
    ]
    add_bullets(s, rf_good, Inches(0.5), Inches(1.85), Inches(5.7), Inches(2.2),
                size=13, marker="  ✓ ", color=ACCENT1, title_color=ACCENT1)

    # Red: stacked model failure
    add_rect(s, Inches(0.4), Inches(4.4), Inches(5.9), Inches(2.65), RGBColor(0x3D, 0x0A, 0x0A))
    add_text(s, "✗  Stacked Ensemble — Failure", Inches(0.5), Inches(4.45),
             Inches(5.7), Inches(0.5), size=15, bold=True, color=ACCENT_R)
    rf_bad = [
        "ROC-AUC 0.5, PR-AUC 0.5, Detection rate 0 — random baseline",
        "Meta-learner trained on data the base models have seen → data leakage",
        "Base model scores on training data not representative of test",
        "Fix: use out-of-fold cross-validation predictions for stacking",
        "Fix: hold out a proper meta-training set before fitting base models",
    ]
    add_bullets(s, rf_bad, Inches(0.5), Inches(4.95), Inches(5.7), Inches(2.0),
                size=12, marker="  ✗ ", color=ACCENT_R, title_color=ACCENT_R)

    # Right: future
    future = [
        "##Future improvements",
        "Fix stacking via proper out-of-fold CV predictions",
        "Add GBM / XGBoost to ensemble for score diversity",
        "Dynamic analysis features (PDF JS execution traces)",
        "Train across more PDF sub-types (forms, exploit kits)",
        "##Recall gap (0.8992 → target 0.95+)",
        "Evasive PDFs use encrypted streams & steganography",
        "Add n-gram features on raw byte sequences",
        "Consider CNN on rendered PDF page images",
        "Semi-supervised: leverage unlabeled PDF corpus",
        "##Deployment",
        "Chunked training → handles streaming PDF telemetry",
        "Feature pipeline must be reproducible at inference time",
    ]
    add_bullets(s, future, Inches(6.6), Inches(1.3), Inches(6.4), Inches(5.8),
                size=13, title_color=ACCENT1)
    add_rect(s, Inches(6.4), Inches(1.25), Pt(2), Inches(5.8), ACCENT1)


# ── final slide ────────────────────────────────────────────────────────────────
def slide_final(prs):
    s = blank_slide(prs); solid_bg(s, BG_DARK)
    add_rect(s, 0, Inches(2.7), W, Inches(2.2), RGBColor(0x1F, 0x77, 0xB4))

    combined = [
        ("EMBER2024 Ensemble",  "Test ROC-AUC 0.9976  ·  Chal PR-AUC 0.6284  ·  Det. Rate 97.86%  (rank ensemble)", ACCENT1),
        ("PDF Random Forest",   "ROC-AUC 0.9888  ·  PR-AUC 0.9909  ·  F1 0.9444  ·  Precision 0.9945",            ACCENT3),
        ("Paper Baseline Ref.", "Test ROC-AUC 0.9968  ·  Chal PR-AUC 0.4725  ·  Det. Rate 66.54%",                 GREY),
    ]
    add_text(s, "Combined Results Summary", Inches(0.6), Inches(2.8),
             Inches(12), Inches(0.6), size=24, bold=True, align=PP_ALIGN.CENTER)
    for i, (name, vals, col) in enumerate(combined):
        y = Inches(3.45) + i * Inches(0.55)
        add_text(s, f"{name}:  {vals}", Inches(0.8), y, Inches(11.7), Inches(0.5),
                 size=13, color=col, align=PP_ALIGN.CENTER)

    add_text(s, "Thank You", Inches(0.6), Inches(5.3), Inches(12), Inches(0.8),
             size=36, bold=True, align=PP_ALIGN.CENTER, color=WHITE)
    add_text(s, "Questions & Discussion", Inches(0.6), Inches(6.0), Inches(12), Inches(0.5),
             size=16, color=GREY, align=PP_ALIGN.CENTER)


# ── assemble ───────────────────────────────────────────────────────────────────
def main():
    prs = new_prs()

    print("Building slides …")
    slide_title(prs)                                                          # 1
    slide_problem(prs)                                                        # 2
    slide_our_approach(prs)                                                   # 3
    slide_our_metrics(prs)                                                    # 4
    slide_our_roc(prs)                                                        # 5
    slide_our_features(prs)                                                   # 6
    slide_our_failure(prs)                                                    # 7
    slide_divider(prs,                                                        # 8
                  "Random Forest — PDF Malware Detection",
                  "Friend's Model  ·  24,000 PDF test samples  ·  Binary Classification",
                  bg=BG_DARK, accent=ACCENT1)
    slide_rf_overview(prs)                                                    # 9
    slide_rf_metrics(prs)                                                     # 10
    slide_rf_curves(prs)                                                      # 11
    slide_rf_features(prs)                                                    # 12
    slide_rf_takeaways(prs)                                                   # 13
    slide_final(prs)                                                          # 14

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"\nSaved → {OUT.resolve()}  ({prs.slides.__len__()} slides)")


if __name__ == "__main__":
    main()
